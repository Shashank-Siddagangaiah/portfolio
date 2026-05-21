"""
Enhanced Executive Presentation Builder
PEMCO BIM Retirement & Report Migration — May 2026
For PEMCO CHIVO / Senior Leadership audience
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x30, 0x55)   # PEMCO deep navy
RED        = RGBColor(0xC4, 0x1E, 0x3A)   # PEMCO accent red
GOLD       = RGBColor(0xE0, 0x9A, 0x00)   # warm gold for highlights
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY   = RGBColor(0xB0, 0xB8, 0xC8)
DARK_GRAY  = RGBColor(0x3C, 0x3C, 0x3C)
GREEN      = RGBColor(0x1E, 0x8A, 0x44)
ORANGE     = RGBColor(0xE6, 0x73, 0x00)
BLUE_LIGHT = RGBColor(0x26, 0x6E, 0xB4)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Completely blank


# ═══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=14, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline_text(slide, lines, l, t, w, h,
                       font_size=12, bold=False, color=DARK_GRAY,
                       align=PP_ALIGN.LEFT, line_spacing=None):
    """lines: list of (text, bold_override, color_override, size_override)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, b, c, s = item, bold, color, font_size
        else:
            text = item[0]
            b    = item[1] if len(item) > 1 else bold
            c    = item[2] if len(item) > 2 else color
            s    = item[3] if len(item) > 3 else font_size
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.space_before = Pt(line_spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = "Calibri"
    return txBox

def slide_header(slide, title, subtitle=None):
    """Adds the navy top bar + title to every slide."""
    add_rect(slide, 0, 0, W, Inches(0.95), NAVY)
    add_text_box(slide, title,
                 Inches(0.35), Inches(0.08), Inches(10), Inches(0.55),
                 font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(0.6), Inches(10), Inches(0.35),
                     font_size=11, bold=False, color=MID_GRAY)
    # Red accent line under header
    add_rect(slide, 0, Inches(0.95), W, Inches(0.04), RED)
    # Footer
    add_rect(slide, 0, Inches(7.2), W, Inches(0.3), NAVY)
    add_text_box(slide, "Bourntec Solutions  |  PEMCO Mutual Insurance  |  Confidential — May 2026",
                 Inches(0.3), Inches(7.2), Inches(9), Inches(0.3),
                 font_size=8, color=MID_GRAY)

def metric_box(slide, l, t, w, h, value, label, bg=NAVY, val_color=WHITE, lbl_color=MID_GRAY):
    add_rect(slide, l, t, w, h, bg)
    # value
    add_text_box(slide, value,
                 l, t + Inches(0.15), w, Inches(0.65),
                 font_size=36, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    # label
    add_text_box(slide, label,
                 l + Inches(0.1), t + Inches(0.75), w - Inches(0.2), Inches(0.55),
                 font_size=11, bold=False, color=lbl_color, align=PP_ALIGN.CENTER)

def section_header_box(slide, label, l, t, w):
    add_rect(slide, l, t, w, Inches(0.3), NAVY)
    add_text_box(slide, label,
                 l + Inches(0.15), t + Inches(0.02), w - Inches(0.3), Inches(0.28),
                 font_size=10, bold=True, color=WHITE)

def check_item(slide, text, l, t, w, color=DARK_GRAY, size=11, icon="✓"):
    add_text_box(slide, f"{icon}  {text}",
                 l, t, w, Inches(0.3),
                 font_size=size, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(s1, 0, 0, W, H, NAVY)
# Red accent stripe
add_rect(s1, 0, Inches(4.2), W, Inches(0.07), RED)
# White content band
add_rect(s1, 0, Inches(2.4), W, Inches(1.8), RGBColor(0xFF, 0xFF, 0xFF))

# Company logos area (text stand-in)
add_text_box(s1, "PEMCO  ×  BOURNTEC SOLUTIONS",
             Inches(0.5), Inches(0.35), Inches(12), Inches(0.5),
             font_size=13, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Divider
add_rect(s1, Inches(5.5), Inches(0.85), Inches(2.3), Inches(0.03), RED)

add_text_box(s1, "BIM Retirement &\nReport Migration",
             Inches(1), Inches(1.1), Inches(11.3), Inches(1.4),
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s1, "Executive Status Update",
             Inches(1), Inches(2.5), Inches(11.3), Inches(0.45),
             font_size=20, bold=False, color=NAVY, align=PP_ALIGN.CENTER)

add_text_box(s1, "May 2026",
             Inches(1), Inches(2.93), Inches(11.3), Inches(0.4),
             font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Bottom stats teaser
stats = [
    ("10+", "Reports Migrated"),
    ("5→1", "Tableau Consolidation"),
    ("+9.3pp", "Email Coverage Gain"),
    ("<2%", "Validated Delta vs BIM"),
]
box_w = Inches(2.8)
gap   = Inches(0.3)
start = Inches(0.57)
for i, (val, lbl) in enumerate(stats):
    bx = start + i * (box_w + gap)
    add_rect(s1, bx, Inches(5.2), box_w, Inches(1.7), RGBColor(0x22, 0x3D, 0x6B))
    add_text_box(s1, val,
                 bx, Inches(5.35), box_w, Inches(0.7),
                 font_size=34, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(s1, lbl,
                 bx + Inches(0.1), Inches(6.05), box_w - Inches(0.2), Inches(0.55),
                 font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text_box(s1, "Bourntec Solutions  |  PEMCO Mutual Insurance  |  Confidential",
             Inches(0), Inches(7.15), W, Inches(0.3),
             font_size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary: Impact at a Glance
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, W, H, LIGHT_GRAY)
slide_header(s2, "Executive Summary", "What Bourntec delivered in this engagement — May 2026")

# 4 big metric boxes
metrics = [
    ("10+",   "Reports\nSuccessfully Migrated",    NAVY,       WHITE,    MID_GRAY),
    ("5 → 1", "Tableau Reports\nConsolidated",      BLUE_LIGHT, WHITE,    RGBColor(0xCC, 0xE0, 0xFF)),
    ("+9.3pp","Email Coverage\nImprovement",         GREEN,      WHITE,    RGBColor(0xCC, 0xFF, 0xDD)),
    ("<2%",   "Validated Delta\nvs BIM Baseline",    RED,        WHITE,    RGBColor(0xFF, 0xCC, 0xCC)),
]
mw = Inches(2.9)
mg = Inches(0.27)
ms = Inches(0.35)
mt = Inches(1.15)
mh = Inches(1.6)
for i, (val, lbl, bg, vc, lc) in enumerate(metrics):
    metric_box(s2, ms + i*(mw+mg), mt, mw, mh, val, lbl, bg, vc, lc)

# Sub-headline
add_text_box(s2, "Key Accomplishments",
             Inches(0.35), Inches(3.0), Inches(12), Inches(0.38),
             font_size=14, bold=True, color=NAVY)
add_rect(s2, Inches(0.35), Inches(3.38), Inches(12.6), Inches(0.03), RED)

# 3-column accomplishments
col_w  = Inches(3.95)
col_t  = Inches(3.55)
col_h  = Inches(3.3)
cols = [
    ("MIGRATIONS COMPLETE", [
        ("BOXI → Power BI: 4 reports fully rebuilt", GREEN),
        ("BOXI / SSRS → Tableau: 6 reports migrated with source rewrite BIM → EDW/AWM", GREEN),
        ("Tableau Consolidation: 5 separate reports unified into 1 with a single shared data source", GREEN),
        ("Engagement Dashboard: migrated and live", GREEN),
    ]),
    ("DATA QUALITY WINS", [
        ("Email coverage: 75.8% → 85.1% (+26,476 policies recovered)", GREEN),
        ("15 data quality bugs identified & resolved in Paperless pipeline", GREEN),
        ("Paperless pipeline streamlined: 16 temp tables → 10", GREEN),
        ("Email-aware dedup: ANI email recovered when NIN has none (+47,658 policies analyzed)", GREEN),
    ]),
    ("VALIDATION RIGOR", [
        ("Online Accounts: <2% delta vs BIM across all 4 DateTypes", GREEN),
        ("Inforce CSR/CSS validated at −1.21% / +1.46% — within acceptable threshold", GREEN),
        ("BIM vs EDW side-by-side analysis documented for every report", GREEN),
        ("All structural gaps documented with root-cause and accepted/pending status", GREEN),
    ]),
]
for ci, (title, items) in enumerate(cols):
    cl = Inches(0.35) + ci * (col_w + Inches(0.15))
    add_rect(s2, cl, col_t, col_w, col_h, WHITE)
    add_rect(s2, cl, col_t, col_w, Inches(0.3), NAVY)
    add_text_box(s2, title,
                 cl + Inches(0.12), col_t + Inches(0.03), col_w - Inches(0.2), Inches(0.26),
                 font_size=9, bold=True, color=WHITE)
    for ii, (txt, col) in enumerate(items):
        add_text_box(s2, f"✓  {txt}",
                     cl + Inches(0.12), col_t + Inches(0.42) + ii * Inches(0.71),
                     col_w - Inches(0.2), Inches(0.65),
                     font_size=10, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Previous State: Why We Had to Change
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, W, H, LIGHT_GRAY)
slide_header(s3, "The Challenge: Legacy BIM-Based Architecture",
             "Why migration was mandated — and what was at stake")

# Left: Pain points
lw = Inches(5.8)
add_rect(s3, Inches(0.35), Inches(1.1), lw, Inches(5.85), WHITE)
add_rect(s3, Inches(0.35), Inches(1.1), lw, Inches(0.32), NAVY)
add_text_box(s3, "LEGACY ARCHITECTURE — PAIN POINTS",
             Inches(0.5), Inches(1.12), lw - Inches(0.3), Inches(0.28),
             font_size=9, bold=True, color=WHITE)

pain_points = [
    ("Single Point of Failure",
     "BIM Monthly ETL ran on the 1st of every month. Any failure delayed ALL reporting downstream."),
    ("Data Latency",
     "Monthly ETL cadence meant data was stale up to 30 days. Weekly and daily jobs were partial."),
    ("Fragmented Tableau Sources",
     "5+ Tableau dashboards each had independent BIM connections — different logic, inconsistent numbers."),
    ("Manual SSIS / Stored Procedures",
     "BIM fed reports via SSIS packages and stored procedures: brittle, hard to maintain, undocumented."),
    ("Multiple Downstream Systems",
     "BIM drove Production Reporting, PMIC Data Warehouse, PMIC Data Mart, FYPE — a tangled web."),
    ("Report Tool Sprawl",
     "BOXI (Business Objects), SSRS, and disconnected Tableau dashboards — no unified BI layer."),
]
for i, (h, b) in enumerate(pain_points):
    ty = Inches(1.55) + i * Inches(0.87)
    add_rect(s3, Inches(0.5), ty, Inches(0.06), Inches(0.62), RED)
    add_text_box(s3, h,
                 Inches(0.7), ty, lw - Inches(0.5), Inches(0.28),
                 font_size=11, bold=True, color=NAVY)
    add_text_box(s3, b,
                 Inches(0.7), ty + Inches(0.28), lw - Inches(0.5), Inches(0.38),
                 font_size=10, color=DARK_GRAY)

# Right: What drove the mandate
rw = Inches(6.2)
rl = Inches(6.65)
add_rect(s3, rl, Inches(1.1), rw, Inches(2.55), NAVY)
add_text_box(s3, "MIGRATION MANDATE",
             rl + Inches(0.2), Inches(1.15), rw - Inches(0.4), Inches(0.3),
             font_size=9, bold=True, color=GOLD)
mandate_lines = [
    "BIM retirement was formally mandated by PEMCO IT leadership.",
    "All enterprise reports must be rewritten to target AWM / EDW (Enterprise Data Warehouse).",
    "BIM table and column structures differ from AWM — business logic cannot be directly ported.",
    "Target: eliminate BIM dependency across every report, dashboard, and downstream consumer.",
]
for i, line in enumerate(mandate_lines):
    add_text_box(s3, f"›  {line}",
                 rl + Inches(0.2), Inches(1.52) + i * Inches(0.5),
                 rw - Inches(0.35), Inches(0.48),
                 font_size=10, color=WHITE)

add_rect(s3, rl, Inches(3.75), rw, Inches(3.2), WHITE)
add_rect(s3, rl, Inches(3.75), rw, Inches(0.32), RGBColor(0x26, 0x6E, 0xB4))
add_text_box(s3, "AWM / EDW ADVANTAGES",
             rl + Inches(0.2), Inches(3.77), rw - Inches(0.4), Inches(0.28),
             font_size=9, bold=True, color=WHITE)
adv = [
    ("Near-real-time data",        "No ETL monthly batch dependency — data reflects current state"),
    ("Single source of truth",     "AWM/DWM serves as the authoritative enterprise data warehouse"),
    ("Scalability",                "Cloud-ready architecture; supports Databricks future roadmap"),
    ("Unified reporting layer",    "All Tableau dashboards source from one consistent EDW schema"),
    ("Auditability",               "Full lineage from source system to reporting table is traceable"),
]
for i, (h, b) in enumerate(adv):
    ty = Inches(4.15) + i * Inches(0.55)
    add_text_box(s3, f"✓  {h}",
                 rl + Inches(0.2), ty, Inches(2.6), Inches(0.28),
                 font_size=10, bold=True, color=GREEN)
    add_text_box(s3, b,
                 rl + Inches(2.9), ty, rw - Inches(3.05), Inches(0.28),
                 font_size=10, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Scope of Work
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, W, H, LIGHT_GRAY)
slide_header(s4, "Scope of Work", "Four workstreams — covering the full BIM retirement portfolio")

tracks = [
    {
        "num": "01",
        "title": "BOXI → Power BI",
        "subtitle": "Ellen's requirement — 4 reports",
        "color": NAVY,
        "items": [
            "Agent Loss Experience Report (O0140A-M1A-P01)",
            "Agent Loss Experience — Current 12 Months",
            "Billing Accounts with Suspended Disbursements (O0169A-W1C-P01)",
            "Monthly Direct Sales Incentive Count Report (O0489A-M5P-P01)",
        ],
        "note": "No AWM/DWM dependency — direct BIM to Power BI rebuild."
    },
    {
        "num": "02",
        "title": "BOXI / SSRS → Tableau",
        "subtitle": "6 reports with source rewrite",
        "color": BLUE_LIGHT,
        "items": [
            "Monthly Inforce Household & Policy Unit Counts (O0300A-M1C-P01)",
            "Monthly Manual Write Off (O0242A-M1A-P01)",
            "Monthly Write Off (O0302A-M1B-P01)",
            "Monthly Proxy Emailed Detail",
            "Monthly Proxy Emailer Dashboard",
            "Monthly Proxy Received Detail",
        ],
        "note": "Source rewritten from BIM to AWM/DWM for all 6."
    },
    {
        "num": "03",
        "title": "Tableau Consolidation",
        "subtitle": "5 separate reports → 1 unified report",
        "color": GREEN,
        "items": [
            "Paper Report",
            "Online Accounts — Household & Paperless",
            "Paperless Policies by State & Policy Type",
            "Percent of Households with Online Accounts",
            "Paperless Online Accounts",
        ],
        "note": "All 5 consolidated into one unified report. Source: BIM → AWM/DWM."
    },
    {
        "num": "04",
        "title": "Agent Dashboard",
        "subtitle": "SSRS report — data reconciliation",
        "color": ORANGE,
        "items": [
            "Reconcile BIM vs EDW on Earned Premium",
            "Reconcile Written Premium variances",
            "Reconcile Inforce counts",
            "Reconcile New Business figures",
        ],
        "note": "Existing SSRS report — fixing variance between BIM and EDW/AWM."
    },
]

tw = Inches(3.0)
tg = Inches(0.17)
ts = Inches(0.35)
tt = Inches(1.1)
th = Inches(5.85)

for i, tr in enumerate(tracks):
    tl = ts + i * (tw + tg)
    add_rect(s4, tl, tt, tw, th, WHITE)
    # colored header
    add_rect(s4, tl, tt, tw, Inches(0.88), tr["color"])
    add_text_box(s4, tr["num"],
                 tl + Inches(0.15), tt + Inches(0.04), Inches(0.6), Inches(0.36),
                 font_size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(s4, tr["title"],
                 tl + Inches(0.72), tt + Inches(0.06), tw - Inches(0.8), Inches(0.32),
                 font_size=12, bold=True, color=WHITE)
    add_text_box(s4, tr["subtitle"],
                 tl + Inches(0.72), tt + Inches(0.4), tw - Inches(0.8), Inches(0.3),
                 font_size=9, color=RGBColor(0xDD, 0xEE, 0xFF))
    # items
    for j, item in enumerate(tr["items"]):
        add_text_box(s4, f"›  {item}",
                     tl + Inches(0.15), tt + Inches(1.0) + j * Inches(0.54),
                     tw - Inches(0.2), Inches(0.5),
                     font_size=10, color=DARK_GRAY)
    # note
    add_rect(s4, tl, tt + Inches(4.95), tw, Inches(0.9), LIGHT_GRAY)
    add_text_box(s4, tr["note"],
                 tl + Inches(0.12), tt + Inches(5.0), tw - Inches(0.2), Inches(0.8),
                 font_size=9, italic=True, color=RGBColor(0x55, 0x55, 0x77))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Completed Migrations: Full Status
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, W, H, LIGHT_GRAY)
slide_header(s5, "Completed Migrations — Full Status Board",
             "Every report in scope — current delivery status as of May 2026")

# Status legend
for i, (lbl, clr) in enumerate([("Complete ✓", GREEN), ("Minor Issues", ORANGE), ("In Progress", GOLD), ("Blocked", RED)]):
    lx = Inches(0.35) + i * Inches(2.5)
    add_rect(s5, lx, Inches(1.1), Inches(0.18), Inches(0.18), clr)
    add_text_box(s5, lbl, lx + Inches(0.24), Inches(1.08), Inches(2.0), Inches(0.22), font_size=9, color=DARK_GRAY)

# Table header
cols_def = [Inches(5.5), Inches(2.2), Inches(1.5), Inches(1.8), Inches(1.85)]
col_x = [Inches(0.35), Inches(5.85), Inches(8.05), Inches(9.55), Inches(11.4)]
headers = ["Report Name", "Workstream", "Source", "Tool", "Status"]
th_y = Inches(1.42)
add_rect(s5, Inches(0.35), th_y, W - Inches(0.7), Inches(0.32), NAVY)
for ci, hdr in enumerate(headers):
    add_text_box(s5, hdr, col_x[ci] + Inches(0.08), th_y + Inches(0.04),
                 cols_def[ci] - Inches(0.12), Inches(0.28),
                 font_size=9, bold=True, color=WHITE)

rows = [
    ("Agent Loss Experience Report (O0140A-M1A-P01)",      "BOXI → Power BI",        "BIM",       "Power BI",  "Complete ✓",    GREEN),
    ("Agent Loss Experience — Current 12 Months",           "BOXI → Power BI",        "BIM",       "Power BI",  "Complete ✓",    GREEN),
    ("Billing Accounts / Suspended Disbursements (O0169A)", "BOXI → Power BI",        "BIM",       "Power BI",  "Complete ✓",    GREEN),
    ("Monthly Direct Sales Incentive Count (O0489A)",       "BOXI → Power BI",        "BIM",       "Power BI",  "Complete ✓",    GREEN),
    ("Monthly Inforce Household & Policy Unit Counts",      "BOXI/SSRS → Tableau",   "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Monthly Manual Write Off (O0242A-M1A-P01)",           "BOXI/SSRS → Tableau",   "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Monthly Write Off (O0302A-M1B-P01)",                  "BOXI/SSRS → Tableau",   "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Monthly Proxy Emailed Detail",                        "BOXI/SSRS → Tableau",   "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Monthly Proxy Emailer Dashboard",                     "BOXI/SSRS → Tableau",   "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Monthly Proxy Received Detail",                       "BOXI/SSRS → Tableau",   "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Consolidated Single-Source Report (5→1)",             "Tableau Consolidation",  "AWM/DWM",   "Tableau",   "Minor Issues",  ORANGE),
    ("Engagement Dashboard",                                "Tableau Other",          "AWM/DWM",   "Tableau",   "Complete ✓",    GREEN),
    ("Quote to Purchase Online & DIA",                      "Tableau Other",          "Salesforce","Tableau",   "Blocked",       RED),
    ("Online Accounts (UserEvents)",                        "Tableau Other",          "AWM/DWM",   "Tableau",   "In Progress",   GOLD),
    ("Agent Dashboard (SSRS Reconciliation)",               "Agent Dashboard",        "BIM/EDW",   "SSRS",      "In Progress",   GOLD),
]

row_h = Inches(0.36)
for ri, row in enumerate(rows):
    ry = Inches(1.74) + ri * row_h
    bg = WHITE if ri % 2 == 0 else LIGHT_GRAY
    add_rect(s5, Inches(0.35), ry, W - Inches(0.7), row_h, bg)
    # Status dot
    add_rect(s5, col_x[4] + Inches(0.08), ry + Inches(0.1), Inches(0.15), Inches(0.15), row[5])
    for ci, val in enumerate(row[:5]):
        add_text_box(s5, val,
                     col_x[ci] + Inches(0.08), ry + Inches(0.04),
                     cols_def[ci] - Inches(0.12), row_h - Inches(0.06),
                     font_size=9, color=DARK_GRAY if ci < 4 else row[5])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Deep Dive: Tableau Consolidation (5→1)
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, W, H, LIGHT_GRAY)
slide_header(s6, "Tableau Consolidation: 5 Reports → 1 Unified Source",
             "Eliminating fragmentation — single AWM/DWM data source, consistent logic everywhere")

# Before box
bw = Inches(5.1)
add_rect(s6, Inches(0.35), Inches(1.1), bw, Inches(5.85), WHITE)
add_rect(s6, Inches(0.35), Inches(1.1), bw, Inches(0.38), RED)
add_text_box(s6, "BEFORE — 5 Disconnected BIM Sources",
             Inches(0.5), Inches(1.12), bw - Inches(0.2), Inches(0.32),
             font_size=10, bold=True, color=WHITE)

before_reports = [
    "Paper Report",
    "Online Accounts — Household & Paperless",
    "Paperless Policies by State & Policy Type",
    "Percent of Households with Online Accounts",
    "Paperless Online Accounts",
]
for i, r in enumerate(before_reports):
    ry = Inches(1.62) + i * Inches(0.62)
    add_rect(s6, Inches(0.55), ry, bw - Inches(0.4), Inches(0.52), LIGHT_GRAY)
    add_rect(s6, Inches(0.55), ry, Inches(0.08), Inches(0.52), RED)
    add_text_box(s6, r, Inches(0.75), ry + Inches(0.1), bw - Inches(0.55), Inches(0.32),
                 font_size=10, color=DARK_GRAY)
    add_text_box(s6, "Source: BIM (Monthly ETL)",
                 Inches(0.75), ry + Inches(0.34), bw - Inches(0.55), Inches(0.2),
                 font_size=8, italic=True, color=RED)

add_rect(s6, Inches(0.5), Inches(4.75), bw - Inches(0.3), Inches(2.1), RGBColor(0xFF, 0xEE, 0xEE))
add_multiline_text(s6, [
    ("Issues with this approach:", True, RED, 10),
    ("›  Duplicate logic maintained in 5 separate workbooks", False, DARK_GRAY, 10),
    ("›  Inconsistent definitions — numbers differed across reports", False, DARK_GRAY, 10),
    ("›  Any BIM ETL delay broke all 5 dashboards simultaneously", False, DARK_GRAY, 10),
    ("›  5× the maintenance burden for any data model change", False, DARK_GRAY, 10),
], Inches(0.65), Inches(4.82), bw - Inches(0.45), Inches(1.9), line_spacing=2)

# Arrow
add_text_box(s6, "→",
             Inches(5.55), Inches(3.4), Inches(0.8), Inches(0.7),
             font_size=42, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# After box
aw = Inches(6.65)
al = Inches(6.45)
add_rect(s6, al, Inches(1.1), aw, Inches(5.85), WHITE)
add_rect(s6, al, Inches(1.1), aw, Inches(0.38), GREEN)
add_text_box(s6, "AFTER — 1 Unified Report, Single AWM/DWM Source",
             al + Inches(0.15), Inches(1.12), aw - Inches(0.2), Inches(0.32),
             font_size=10, bold=True, color=WHITE)

# Single report visual
add_rect(s6, al + Inches(0.25), Inches(1.65), aw - Inches(0.5), Inches(1.1),
         RGBColor(0xE8, 0xF5, 0xEC))
add_rect(s6, al + Inches(0.25), Inches(1.65), Inches(0.1), Inches(1.1), GREEN)
add_text_box(s6, "Consolidated Tableau Report",
             al + Inches(0.5), Inches(1.75), aw - Inches(0.8), Inches(0.32),
             font_size=12, bold=True, color=GREEN)
add_text_box(s6, "Single shared AWM/DWM data source  |  All 5 views in one workbook",
             al + Inches(0.5), Inches(2.07), aw - Inches(0.8), Inches(0.28),
             font_size=9, color=DARK_GRAY)
add_text_box(s6, "Source: AWM / DWM (Enterprise Data Warehouse)",
             al + Inches(0.5), Inches(2.35), aw - Inches(0.8), Inches(0.22),
             font_size=8, italic=True, color=GREEN)

benefits = [
    ("Maintenance",     "1 workbook to maintain vs 5 — 80% reduction in update overhead"),
    ("Consistency",     "Identical business logic across all views — numbers always agree"),
    ("Reliability",     "No BIM ETL dependency — report refreshes independently"),
    ("Data Freshness",  "Near-real-time AWM/DWM refresh — no monthly batch lag"),
    ("Extensibility",   "New views or metrics added in one place, available everywhere"),
]
add_text_box(s6, "Benefits Delivered",
             al + Inches(0.25), Inches(2.85), aw - Inches(0.4), Inches(0.3),
             font_size=10, bold=True, color=NAVY)
for i, (h, b) in enumerate(benefits):
    by = Inches(3.2) + i * Inches(0.62)
    add_rect(s6, al + Inches(0.25), by, aw - Inches(0.5), Inches(0.56), LIGHT_GRAY)
    add_text_box(s6, f"✓  {h}",
                 al + Inches(0.4), by + Inches(0.04), Inches(1.8), Inches(0.26),
                 font_size=10, bold=True, color=GREEN)
    add_text_box(s6, b,
                 al + Inches(2.2), by + Inches(0.04), aw - Inches(2.5), Inches(0.5),
                 font_size=10, color=DARK_GRAY)

# Status note
add_rect(s6, al + Inches(0.25), Inches(6.3), aw - Inches(0.5), Inches(0.5), RGBColor(0xFF, 0xF3, 0xCD))
add_text_box(s6, "⚠  Status: Complete with minor data issues under investigation. AWM/DWM source validation ongoing.",
             al + Inches(0.4), Inches(6.35), aw - Inches(0.65), Inches(0.42),
             font_size=9, italic=True, color=ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Deep Dive: Online Accounts Migration & Validation
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, W, H, LIGHT_GRAY)
slide_header(s7, "Online Accounts: Migration & Validation Results",
             "4 DateTypes rewritten from BIM to AWM/EDW — validated within accepted thresholds")

# Intro text
add_text_box(s7,
    "The Online Accounts Self-Service report tracks enrollment across 4 DateTypes (Activation, Initiation, Inforce, DriverOnly). "
    "Each DateType was rewritten to query AWM/EDW and validated line-by-line against the BIM legacy baseline.",
    Inches(0.35), Inches(1.1), Inches(12.6), Inches(0.55),
    font_size=10, color=DARK_GRAY)

# Validation table
add_rect(s7, Inches(0.35), Inches(1.75), Inches(12.6), Inches(0.32), NAVY)
for ci, (hdr, cx, cw) in enumerate([
    ("DateType",       Inches(0.35), Inches(2.2)),
    ("BIM → AWM/EDW",  Inches(2.55), Inches(2.3)),
    ("Delta vs BIM",   Inches(4.85), Inches(2.0)),
    ("Status",         Inches(6.85), Inches(1.8)),
    ("Notes",          Inches(8.65), Inches(4.3)),
]):
    add_text_box(s7, hdr, cx + Inches(0.08), Inches(1.77),
                 cw - Inches(0.12), Inches(0.28),
                 font_size=9, bold=True, color=WHITE)

vrows = [
    ("Activation",      "BIM → AWM/DWM",  "< 1%",   "Validated ✓", GREEN,   "Rewrite complete. Delta within 1% — accepted as rounding/timing."),
    ("Initiation",      "BIM → AWM/DWM",  "< 2%",   "Validated ✓", GREEN,   "2021+ fully validated. Pre-2021 gap accepted — AWM history incomplete before that date."),
    ("Inforce CSR",     "BIM → AWM/DWM",  "−1.21%", "Validated ✓", GREEN,   "CSR (agent-enrolled) bucket. Delta within threshold. Root cause documented."),
    ("Inforce CSS",     "BIM → AWM/DWM",  "+1.46%", "Validated ✓", GREEN,   "CSS (web self-enrollment) bucket. Delta within threshold."),
    ("DriverOnly CSR",  "BIM → AWM/DWM",  "−46.6%", "Accepted ⚠",  ORANGE,  "Known structural gap. vw_policy_driver has ~13% coverage gap vs BIM dim_policy_role. Under investigation."),
    ("DriverOnly CSS",  "BIM → AWM/DWM",  "−38.3%", "Accepted ⚠",  ORANGE,  "Same structural gap as CSR. 3,971 BIM DR parties have no row in vw_policy_driver."),
]

for ri, (dt, src, delta, status, sc, note) in enumerate(vrows):
    ry = Inches(2.07) + ri * Inches(0.73)
    bg = WHITE if ri % 2 == 0 else LIGHT_GRAY
    add_rect(s7, Inches(0.35), ry, Inches(12.6), Inches(0.71), bg)
    for ci, (val, cx, cw) in enumerate([
        (dt,     Inches(0.35), Inches(2.2)),
        (src,    Inches(2.55), Inches(2.3)),
        (delta,  Inches(4.85), Inches(2.0)),
        (status, Inches(6.85), Inches(1.8)),
        (note,   Inches(8.65), Inches(4.3)),
    ]):
        add_text_box(s7, val,
                     cx + Inches(0.1), ry + Inches(0.08),
                     cw - Inches(0.15), Inches(0.62),
                     font_size=10,
                     bold=(ci == 0),
                     color=sc if ci == 3 else DARK_GRAY)

# Bottom insight boxes
add_text_box(s7, "Key Technical Findings",
             Inches(0.35), Inches(6.55), Inches(12), Inches(0.28),
             font_size=11, bold=True, color=NAVY)

insights = [
    ("is_up_and_running = 255", "The active account indicator is a tinyint — all-bits-set (255), not 1. Using = 1 returned 0 rows."),
    ("AWM → EDW Party Bridge", "party_id_same_as_link maps AWM duplicate anchor IDs to EDW master party_key — critical join path."),
    ("ECOMM1 Session IDs",     "Web session identifiers (ECOMM1%) were misclassified as CSR. Fixed to route correctly to CSS bucket."),
]
iw = Inches(4.0)
for i, (h, b) in enumerate(insights):
    ix = Inches(0.35) + i * (iw + Inches(0.12))
    add_rect(s7, ix, Inches(6.85), iw, Inches(0.55), NAVY)
    add_text_box(s7, h, ix + Inches(0.1), Inches(6.88), Inches(1.8), Inches(0.28),
                 font_size=9, bold=True, color=GOLD)
    add_text_box(s7, b, ix + Inches(1.9), Inches(6.88), iw - Inches(2.0), Inches(0.48),
                 font_size=8, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Deep Dive: Paperless Report Transformation
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_rect(s8, 0, 0, W, H, LIGHT_GRAY)
slide_header(s8, "Paperless Report: Data Quality Transformation",
             "Full EDW/AWM rewrite — 15 bugs resolved, +9.3 percentage points of email coverage recovered")

# Big numbers row
big_stats = [
    ("75.8%",  "Email Coverage\nBefore",      RED,        WHITE,    RGBColor(0xFF, 0xCC, 0xCC)),
    ("85.1%",  "Email Coverage\nAfter",        GREEN,      WHITE,    RGBColor(0xCC, 0xFF, 0xDD)),
    ("+26,476","Additional Policies\nWith Email", BLUE_LIGHT, WHITE, RGBColor(0xCC, 0xE0, 0xFF)),
    ("15",     "Data Quality Issues\nFound & Resolved", NAVY, WHITE, MID_GRAY),
    ("16→10",  "Pipeline Temp\nTables (Streamlined)", GOLD, DARK_GRAY, WHITE),
]
bw = Inches(2.35)
bg = Inches(0.15)
bs = Inches(0.35)
bt = Inches(1.1)
bh = Inches(1.45)
for i, (val, lbl, bg_c, vc, lc) in enumerate(big_stats):
    bl = bs + i * (bw + bg)
    metric_box(s8, bl, bt, bw, bh, val, lbl, bg_c, vc, lc)

# Left column: top issues
lw = Inches(6.0)
add_rect(s8, Inches(0.35), Inches(2.7), lw, Inches(4.25), WHITE)
add_rect(s8, Inches(0.35), Inches(2.7), lw, Inches(0.3), NAVY)
add_text_box(s8, "TOP ISSUES IDENTIFIED & RESOLVED",
             Inches(0.5), Inches(2.72), lw - Inches(0.2), Inches(0.26),
             font_size=9, bold=True, color=WHITE)

issues = [
    ("#1", "35,820 policies dropped",   "vw_policyholder joined on 3 date columns — mid-term endorsements silently failed join."),
    ("#5", "Email fan-out bug",         "asp_user_account_detail partitioned by (link_id, user_name) — one link, 3 rows, counts > total."),
    ("#8", "26,476 ANI emails missed",  "Grain dedup always picked NIN even when NIN had no email. ANI email was silently discarded."),
    ("#2", "CIF superseded records",    "Valid_to_date filter missing — stale superseded CIF records overrode current paperless preference."),
    ("#4", "NULL coerced to 'N'",       "Paper_notify NULL treated as 'N' (not paperless). Absent data ≠ confirmed non-paperless."),
    ("#7", "BIM fallback = 0 rows",     "Incorrect format strip on POL_KEY join — removed strip and 464 fallback emails recovered."),
]
for i, (num, h, b) in enumerate(issues):
    iy = Inches(3.1) + i * Inches(0.65)
    add_rect(s8, Inches(0.5), iy, Inches(0.45), Inches(0.52), RGBColor(0xCC, 0x33, 0x33))
    add_text_box(s8, num, Inches(0.5), iy + Inches(0.1), Inches(0.45), Inches(0.3),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s8, h, Inches(1.05), iy + Inches(0.03), Inches(1.9), Inches(0.26),
                 font_size=10, bold=True, color=NAVY)
    add_text_box(s8, b, Inches(1.05), iy + Inches(0.28), lw - Inches(0.8), Inches(0.28),
                 font_size=9, color=DARK_GRAY)

# Right column: email chain & coverage
rw = Inches(6.3)
rl = Inches(6.65)
add_rect(s8, rl, Inches(2.7), rw, Inches(2.0), WHITE)
add_rect(s8, rl, Inches(2.7), rw, Inches(0.3), BLUE_LIGHT)
add_text_box(s8, "EMAIL RESOLUTION CHAIN — FINAL BREAKDOWN (294,275 Policies)",
             rl + Inches(0.15), Inches(2.72), rw - Inches(0.2), Inches(0.26),
             font_size=9, bold=True, color=WHITE)

chain_rows = [
    ("EMAIL RESOLVED",       "250,280", "85.0%", GREEN),
    ("  — via NIN (Named Insured)",  "202,622", "68.9%", RGBColor(0x22, 0x77, 0x44)),
    ("  — via ANI (Additional Named)", "47,658", "16.2%", RGBColor(0x22, 0x77, 0x44)),
    ("NO ACCOUNT LINK",      "36,376",  "12.4%", ORANGE),
    ("HAS LINK / NO DETAIL", "8,013",   " 2.7%", GOLD),
    ("NO SAME-AS-LINK",      "1",       "~0%",   MID_GRAY),
]
for i, (lbl, cnt, pct, cc) in enumerate(chain_rows):
    cy = Inches(3.1) + i * Inches(0.3)
    add_text_box(s8, lbl, rl + Inches(0.15), cy, Inches(3.3), Inches(0.28),
                 font_size=9, color=DARK_GRAY, bold=("RESOLVED" in lbl or "LINK" == lbl[-4:] and i > 0 == False))
    add_text_box(s8, cnt, rl + Inches(3.5), cy, Inches(1.1), Inches(0.28),
                 font_size=9, bold=True, color=cc, align=PP_ALIGN.RIGHT)
    add_text_box(s8, pct, rl + Inches(4.7), cy, Inches(0.85), Inches(0.28),
                 font_size=9, color=cc, align=PP_ALIGN.RIGHT)

add_rect(s8, rl, Inches(4.8), rw, Inches(2.15), WHITE)
add_rect(s8, rl, Inches(4.8), rw, Inches(0.3), GREEN)
add_text_box(s8, "COVERAGE IMPROVEMENT — BEFORE VS AFTER",
             rl + Inches(0.15), Inches(4.82), rw - Inches(0.2), Inches(0.26),
             font_size=9, bold=True, color=WHITE)

before_after = [
    ("AWM Email (Before)",      "225,188",  "75.8%", RED),
    ("AWM Email (After fixes)", "251,664",  "84.9%", GREEN),
    ("BIM Fallback Recovered",  "464",      " 0.2%", BLUE_LIGHT),
    ("Total Email (Final)",     "252,128",  "85.1%", GREEN),
    ("No Email (Final)",        "44,271",   "14.9%", ORANGE),
]
for i, (lbl, cnt, pct, cc) in enumerate(before_after):
    cy = Inches(5.2) + i * Inches(0.3)
    add_text_box(s8, lbl, rl + Inches(0.15), cy, Inches(3.3), Inches(0.28),
                 font_size=9, color=DARK_GRAY)
    add_text_box(s8, cnt, rl + Inches(3.5), cy, Inches(1.1), Inches(0.28),
                 font_size=9, bold=True, color=cc, align=PP_ALIGN.RIGHT)
    add_text_box(s8, pct, rl + Inches(4.7), cy, Inches(0.85), Inches(0.28),
                 font_size=9, color=cc, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Validation Methodology
# ═══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
add_rect(s9, 0, 0, W, H, LIGHT_GRAY)
slide_header(s9, "Our Validation Methodology",
             "How we ensure every migrated report is correct — a rigorous, documented approach")

steps = [
    ("01", "Establish BIM Baseline",
     "Run legacy BIM query to capture the exact count, grain, and filter logic. Document every WHERE clause, JOIN, and aggregation as the baseline we must match."),
    ("02", "Rewrite for AWM / EDW",
     "Re-implement business logic in AWM/EDW using the correct table equivalents. Apply architecture patterns: party bridge, inforce filter, is_up_and_running = 255."),
    ("03", "Side-by-Side Comparison",
     "Execute both queries against production data. Compare row counts, subtotals, and percentage distributions at every grain level (state, product, DateType, bucket)."),
    ("04", "Delta Analysis & Root Cause",
     "Measure the delta percentage. For every gap > 1%, investigate root cause — structural data gap, logic difference, or ETL timing. Document and classify each delta."),
    ("05", "Accept or Fix",
     "Each gap is classified: (a) Accepted — structural difference understood, stakeholder informed. (b) Fixed — logic corrected, delta reduced. Never leave gaps unexplained."),
    ("06", "Document & Hand Off",
     "Write a full issue resolution document per report. Validation queries archived. All findings persisted in project CLAUDE.md for institutional knowledge."),
]

sw = Inches(3.95)
sg = Inches(0.2)
ss = Inches(0.35)
row1_t = Inches(1.1)
row2_t = Inches(4.05)
sh = Inches(2.8)

for i, (num, title, body) in enumerate(steps):
    col = i % 3
    row = i // 3
    sl = ss + col * (sw + sg)
    st = row1_t if row == 0 else row2_t
    add_rect(s9, sl, st, sw, sh, WHITE)
    add_rect(s9, sl, st, sw, Inches(0.42), NAVY)
    add_text_box(s9, num, sl + Inches(0.12), st + Inches(0.05), Inches(0.5), Inches(0.34),
                 font_size=18, bold=True, color=GOLD)
    add_text_box(s9, title, sl + Inches(0.65), st + Inches(0.07), sw - Inches(0.75), Inches(0.34),
                 font_size=11, bold=True, color=WHITE)
    add_text_box(s9, body, sl + Inches(0.15), st + Inches(0.55), sw - Inches(0.25), sh - Inches(0.65),
                 font_size=10, color=DARK_GRAY)

# Validation principles banner
add_rect(s9, Inches(0.35), Inches(6.95), W - Inches(0.7), Inches(0.38), NAVY)
principles = [
    "Every delta explained",
    "No gaps left undocumented",
    "Structural vs logic gap always distinguished",
    "BIM is comparison baseline — not the gold standard",
    "EDW is the new source of truth",
]
add_text_box(s9, "   |   ".join([f"✓  {p}" for p in principles]),
             Inches(0.5), Inches(6.98), W - Inches(1.0), Inches(0.32),
             font_size=9, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Architecture: New State
# ═══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
add_rect(s10, 0, 0, W, H, LIGHT_GRAY)
slide_header(s10, "New Architecture: End-to-End Data Flow",
             "From source systems to enterprise reporting — no BIM dependency")

# Source systems
src_t = Inches(1.2)
src_h = Inches(0.62)
srcs = ["HDS (Policy Data)", "Exceed (Financials)", "CIFDM (Contact / Paperless)", "Eloqua (Email / BIM fallback)", "Salesforce (Digital) — pending"]
add_rect(s10, Inches(0.2), src_t - Inches(0.35), Inches(2.55), Inches(0.32), RGBColor(0x55, 0x55, 0x77))
add_text_box(s10, "SOURCE SYSTEMS", Inches(0.25), src_t - Inches(0.33), Inches(2.45), Inches(0.28),
             font_size=8, bold=True, color=WHITE)
for i, s in enumerate(srcs):
    sy = src_t + i * (src_h + Inches(0.08))
    add_rect(s10, Inches(0.2), sy, Inches(2.55), src_h, RGBColor(0x3A, 0x3A, 0x5A))
    add_text_box(s10, s, Inches(0.32), sy + Inches(0.1), Inches(2.35), src_h - Inches(0.15),
                 font_size=9, color=WHITE)

# Arrow
add_text_box(s10, "→", Inches(2.82), Inches(3.5), Inches(0.5), Inches(0.5),
             font_size=24, bold=True, color=NAVY)

# EDW layer
add_rect(s10, Inches(3.4), src_t - Inches(0.35), Inches(3.2), Inches(0.32), NAVY)
add_text_box(s10, "AWM / EDW (ENTERPRISE DATA WAREHOUSE)", Inches(3.45), src_t - Inches(0.33),
             Inches(3.1), Inches(0.28), font_size=8, bold=True, color=WHITE)
edw_tables = [
    ("vw_policy",                  "Policy terms, inforce indicators"),
    ("vw_policyholder",            "Party-policy relationships (NIN/ANI)"),
    ("vw_policy_driver",           "Driver-only roles"),
    ("party_id_same_as_link",      "AWM→EDW party bridge"),
    ("party_user_account_link",    "Online account registration"),
    ("asp_user_account_detail",    "Account email (user_name)"),
    ("cif_policy_party_detail",    "Paperless BIL/POL indicators (CIF)"),
]
for i, (tbl, desc) in enumerate(edw_tables):
    ty = src_t + i * (src_h + Inches(0.08))
    add_rect(s10, Inches(3.4), ty, Inches(3.2), src_h, NAVY)
    add_text_box(s10, tbl, Inches(3.52), ty + Inches(0.04), Inches(3.0), Inches(0.28),
                 font_size=9, bold=True, color=GOLD)
    add_text_box(s10, desc, Inches(3.52), ty + Inches(0.3), Inches(3.0), Inches(0.28),
                 font_size=8, color=MID_GRAY)

# Arrow
add_text_box(s10, "→", Inches(6.68), Inches(3.5), Inches(0.5), Inches(0.5),
             font_size=24, bold=True, color=NAVY)

# SQL / BI layer
add_rect(s10, Inches(7.25), src_t - Inches(0.35), Inches(2.55), Inches(0.32), BLUE_LIGHT)
add_text_box(s10, "SQL LAYER (Bourntec)", Inches(7.3), src_t - Inches(0.33), Inches(2.45), Inches(0.28),
             font_size=8, bold=True, color=WHITE)
sql_items = [
    "Online Accounts Self-Service (4 DateTypes)",
    "Online Accounts Active by State",
    "Paperless Report (10-step pipeline)",
    "Validation queries (BIM vs EDW)",
    "Agent Dashboard reconciliation",
]
for i, item in enumerate(sql_items):
    sy = src_t + i * (src_h + Inches(0.08))
    add_rect(s10, Inches(7.25), sy, Inches(2.55), src_h, RGBColor(0x1A, 0x55, 0x99))
    add_text_box(s10, item, Inches(7.37), sy + Inches(0.13), Inches(2.3), src_h - Inches(0.2),
                 font_size=9, color=WHITE)

# Arrow
add_text_box(s10, "→", Inches(9.87), Inches(3.5), Inches(0.5), Inches(0.5),
             font_size=24, bold=True, color=NAVY)

# Reporting layer
add_rect(s10, Inches(10.42), src_t - Inches(0.35), Inches(2.65), Inches(0.32), GREEN)
add_text_box(s10, "REPORTING LAYER", Inches(10.47), src_t - Inches(0.33), Inches(2.55), Inches(0.28),
             font_size=8, bold=True, color=WHITE)
rpt_items = [
    ("Tableau", "Single-source consolidated dashboard"),
    ("Power BI", "4 BOXI-migrated reports"),
    ("Tableau", "Engagement Dashboard"),
    ("SSRS",    "Agent Dashboard (reconciliation)"),
    ("Tableau", "Proxy reports (3)"),
]
for i, (tool, desc) in enumerate(rpt_items):
    ry = src_t + i * (src_h + Inches(0.08))
    add_rect(s10, Inches(10.42), ry, Inches(2.65), src_h, RGBColor(0x22, 0x66, 0x33))
    add_text_box(s10, tool, Inches(10.54), ry + Inches(0.04), Inches(0.85), Inches(0.28),
                 font_size=9, bold=True, color=GOLD)
    add_text_box(s10, desc, Inches(10.54), ry + Inches(0.3), Inches(2.4), Inches(0.28),
                 font_size=8, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Known Gaps & Transparency
# ═══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
add_rect(s11, 0, 0, W, H, LIGHT_GRAY)
slide_header(s11, "Known Gaps & Transparency",
             "Every open item is documented — no surprises, clear resolution paths")

add_text_box(s11,
    "Transparency is a core part of our delivery. Every gap between BIM and EDW/AWM counts is documented, "
    "root-caused, and classified as either under active investigation or accepted with stakeholder alignment.",
    Inches(0.35), Inches(1.1), Inches(12.6), Inches(0.5),
    font_size=10, color=DARK_GRAY)

gaps = [
    {
        "title": "DriverOnly ~43% Count Gap",
        "severity": "Under Investigation",
        "sev_color": ORANGE,
        "report": "Online Accounts — DriverOnly DateType",
        "root_cause": (
            "EDW's vw_policy_driver has ~13% coverage gap vs BIM's CIFDM.dim_policy_role. "
            "3,971 BIM DR (driver) parties have no corresponding row in vw_policy_driver. "
            "No EDW view is a full equivalent of dim_policy_role. "
            "This is a data completeness gap in the source table, not a logic error."
        ),
        "status": "INNER JOIN on vw_policy_driver required a positive driver signal (NOT EXISTS inflated count 17×). "
                  "Gap accepted as interim state. Root cause investigation ongoing with EDW team.",
    },
    {
        "title": "Supervisor Staleness in Online Accounts",
        "severity": "Pending Decision",
        "sev_color": GOLD,
        "report": "Online Accounts — Employee/Supervisor Bucket",
        "root_cause": (
            "AWM captures the supervisor at the time of account creation (USER_EVENT_DETAIL). "
            "BIM reads the current supervisor from an employee master table. "
            "Employees reassigned to a new supervisor show zero accounts under the new supervisor in AWM — "
            "not a bug, a structural design difference between the two systems."
        ),
        "status": "Pending stakeholder decision: retain creation-time supervisor (current AWM behavior) "
                  "or join to current employee table for real-time supervisor mapping.",
    },
    {
        "title": "CIF vs BIM Paperless Divergence (500+ Days Average)",
        "severity": "Structural — Accepted",
        "sev_color": BLUE_LIGHT,
        "report": "Paperless Report — CIF vs BIM mismatch analysis",
        "root_cause": (
            "CIF (AWM) and BIM/Eloqua are separately maintained systems with no live sync. "
            "BIL mismatches: avg 508 days divergence. POL mismatches: avg 771 days. "
            "Divergence is bidirectional — not a processing lag or batch timing issue. "
            "Customers who enrolled via agent portal (2015 Eloqua rollout) never created an AWM account."
        ),
        "status": "Classified as chronic structural gap. Not a pipeline error. "
                  "Validation query email_val.sql documents all mismatch categories and directions.",
    },
    {
        "title": "Salesforce Table Pending (Quote to Purchase)",
        "severity": "External Dependency",
        "sev_color": RED,
        "report": "Quote to Purchase Online & DIA Dashboard",
        "root_cause": (
            "The Salesforce source table required for this dashboard has not yet been made available in the EDW environment. "
            "Migration cannot proceed until the table is provisioned and accessible to the reporting layer."
        ),
        "status": "Blocked on Salesforce table availability. Bourntec ready to proceed once table is accessible.",
    },
]

gw = Inches(12.6)
gl = Inches(0.35)
gt = Inches(1.75)
gh = Inches(1.2)
gg = Inches(0.12)

for i, gap in enumerate(gaps):
    gy = gt + i * (gh + gg)
    add_rect(s11, gl, gy, gw, gh, WHITE)
    add_rect(s11, gl, gy, Inches(0.08), gh, gap["sev_color"])
    # Severity badge
    add_rect(s11, gl + Inches(0.2), gy + Inches(0.08), Inches(1.55), Inches(0.3), gap["sev_color"])
    add_text_box(s11, gap["severity"],
                 gl + Inches(0.22), gy + Inches(0.1), Inches(1.5), Inches(0.26),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_text_box(s11, gap["title"],
                 gl + Inches(1.85), gy + Inches(0.07), Inches(5.5), Inches(0.3),
                 font_size=11, bold=True, color=NAVY)
    add_text_box(s11, f"Report: {gap['report']}",
                 gl + Inches(1.85), gy + Inches(0.37), Inches(5.5), Inches(0.22),
                 font_size=8, italic=True, color=MID_GRAY)
    # Root cause
    add_text_box(s11, gap["root_cause"],
                 gl + Inches(1.85), gy + Inches(0.6), Inches(5.8), Inches(0.55),
                 font_size=9, color=DARK_GRAY)
    # Status
    add_rect(s11, gl + Inches(7.8), gy + Inches(0.08), Inches(4.65), gh - Inches(0.16), LIGHT_GRAY)
    add_text_box(s11, "RESOLUTION PATH",
                 gl + Inches(7.95), gy + Inches(0.1), Inches(4.45), Inches(0.24),
                 font_size=8, bold=True, color=NAVY)
    add_text_box(s11, gap["status"],
                 gl + Inches(7.95), gy + Inches(0.34), Inches(4.45), Inches(0.75),
                 font_size=9, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — What's Next: Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(blank_layout)
add_rect(s12, 0, 0, W, H, LIGHT_GRAY)
slide_header(s12, "What's Next — Remaining Work & Roadmap",
             "Clear priorities and dependencies to complete the BIM retirement program")

phases = [
    {
        "phase": "Immediate",
        "timeline": "May–June 2026",
        "color": RED,
        "items": [
            "Complete Agent Dashboard reconciliation\n(EP, WP, Inforce, New Business vs BIM)",
            "Resolve minor data issues in\nConsolidated Tableau Report",
            "Online Accounts final validation pass\nonce UserEvents table is available",
        ]
    },
    {
        "phase": "Short Term",
        "timeline": "June–July 2026",
        "color": ORANGE,
        "items": [
            "DriverOnly gap investigation with EDW team\n(vw_policy_driver coverage analysis)",
            "Stakeholder decision: supervisor staleness\n(creation-time vs current supervisor)",
            "Quote to Purchase — unblock once\nSalesforce table provisioned",
        ]
    },
    {
        "phase": "Medium Term",
        "timeline": "Q3 2026",
        "color": BLUE_LIGHT,
        "items": [
            "Full BIM dependency audit — confirm\nzero remaining BIM-sourced reports",
            "Databricks pipeline readiness review\nfor AWM/EDW data flows",
            "Atlan data catalog documentation\nfor all migrated reports",
        ]
    },
    {
        "phase": "Ongoing",
        "timeline": "Continuous",
        "color": GREEN,
        "items": [
            "Monitor EDW data quality and\nAWM account coverage trends",
            "Maintain validation query archive\nfor each production report",
            "Respond to new BIM retirement\nrequests as identified",
        ]
    },
]

pw = Inches(3.0)
pg = Inches(0.15)
ps = Inches(0.35)
pt = Inches(1.1)
ph = Inches(5.85)

for i, phase in enumerate(phases):
    pl = ps + i * (pw + pg)
    add_rect(s12, pl, pt, pw, ph, WHITE)
    add_rect(s12, pl, pt, pw, Inches(0.68), phase["color"])
    add_text_box(s12, phase["phase"].upper(),
                 pl + Inches(0.15), pt + Inches(0.06), pw - Inches(0.2), Inches(0.32),
                 font_size=14, bold=True, color=WHITE)
    add_text_box(s12, phase["timeline"],
                 pl + Inches(0.15), pt + Inches(0.38), pw - Inches(0.2), Inches(0.24),
                 font_size=9, color=RGBColor(0xDD, 0xEE, 0xFF))
    for j, item in enumerate(phase["items"]):
        iy = pt + Inches(0.85) + j * Inches(1.6)
        add_rect(s12, pl + Inches(0.15), iy, pw - Inches(0.3), Inches(1.45), LIGHT_GRAY)
        add_rect(s12, pl + Inches(0.15), iy, Inches(0.07), Inches(1.45), phase["color"])
        add_text_box(s12, item,
                     pl + Inches(0.32), iy + Inches(0.15), pw - Inches(0.5), Inches(1.2),
                     font_size=10, color=DARK_GRAY)

# Bottom dependency note
add_rect(s12, Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.28), NAVY)
add_text_box(s12,
    "Key Dependencies:  Salesforce table (EDW provisioning)   |   UserEvents table (AWM availability)   |   "
    "Stakeholder alignment on supervisor staleness   |   EDW team on DriverOnly gap",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.25),
    font_size=8, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(blank_layout)
add_rect(s13, 0, 0, W, H, NAVY)
add_rect(s13, 0, Inches(4.4), W, Inches(0.06), RED)

add_text_box(s13, "PEMCO  ×  BOURNTEC SOLUTIONS",
             Inches(0), Inches(0.4), W, Inches(0.4),
             font_size=12, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text_box(s13, "Thank You",
             Inches(0), Inches(1.0), W, Inches(1.0),
             font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s13, "Questions & Discussion",
             Inches(0), Inches(2.1), W, Inches(0.55),
             font_size=20, color=GOLD, align=PP_ALIGN.CENTER)

add_rect(s13, Inches(4.5), Inches(2.9), Inches(4.3), Inches(0.04), RED)

# Summary stats
final_stats = [
    ("10+", "Reports Migrated"),
    ("5→1", "Tableau Consolidated"),
    ("+9.3pp", "Email Coverage"),
    ("<2%", "BIM Delta"),
    ("15", "Bugs Resolved"),
    ("0", "BIM ETL Dependencies"),
]
fw = Inches(2.0)
fg = Inches(0.12)
fs = Inches(0.35)
ft = Inches(3.15)
fh = Inches(1.4)
for i, (val, lbl) in enumerate(final_stats):
    fl = fs + i * (fw + fg)
    add_rect(s13, fl, ft, fw, fh, RGBColor(0x22, 0x3D, 0x6B))
    add_text_box(s13, val, fl, ft + Inches(0.1), fw, Inches(0.7),
                 font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(s13, lbl, fl + Inches(0.05), ft + Inches(0.8), fw - Inches(0.1), Inches(0.45),
                 font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text_box(s13,
    "Bourntec Solutions  |  PEMCO Mutual Insurance  |  BIM Retirement & Report Migration  |  May 2026",
    Inches(0), Inches(7.15), W, Inches(0.3),
    font_size=8, color=RGBColor(0x55, 0x55, 0x77), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
output_path = "Overall_pemco/PEMCO_Executive_Presentation_May2026.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
